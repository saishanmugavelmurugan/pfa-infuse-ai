"""
Infuse Pitch Deck Generator
Creates a professional pitch deck presentation with brand colors
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

# Infuse Brand Colors
INFUSE_GOLD = RGBColor(0xFF, 0xDA, 0x7B)      # #FFDA7B - Lightest Gold
INFUSE_ORANGE = RGBColor(0xFF, 0x9A, 0x3B)    # #FF9A3B - Medium Orange
INFUSE_DEEP_ORANGE = RGBColor(0xE5, 0x5A, 0x00)  # #E55A00 - Deep Orange
INFUSE_DARK_ORANGE = RGBColor(0xC6, 0x47, 0x00)  # #C64700 - Darkest Orange
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x1F, 0x29, 0x37)
LIGHT_GRAY = RGBColor(0x6B, 0x72, 0x80)

def add_gradient_background(slide, color1, color2):
    """Add a gradient-like background using shapes"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color1
    background.line.fill.background()
    # Move to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_title_slide(prs, title, subtitle=""):
    """Create a title slide with gradient background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_gradient_background(slide, DARK_GRAY, DARK_GRAY)
    
    # Add accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), Inches(13.33), Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = INFUSE_ORANGE
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.33), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = INFUSE_GOLD
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullet_points, accent_color=INFUSE_ORANGE):
    """Create a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    # Top accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """Create a two-column comparison slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    # Top accent bar - gradient effect
    accent1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(6.665), Inches(0.15))
    accent1.fill.solid()
    accent1.fill.fore_color.rgb = INFUSE_ORANGE
    accent1.line.fill.background()
    
    accent2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.665), 0, Inches(6.665), Inches(0.15))
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = INFUSE_DEEP_ORANGE
    accent2.line.fill.background()
    
    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Left Column
    left_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.2), Inches(6), Inches(0.5))
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = INFUSE_ORANGE
    left_header.line.fill.background()
    
    left_title_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.25), Inches(6), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    left_content = slide.shapes.add_textbox(Inches(0.4), Inches(1.9), Inches(5.8), Inches(5))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✓ {point}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    # Right Column
    right_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.2), Inches(6), Inches(0.5))
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = INFUSE_DEEP_ORANGE
    right_header.line.fill.background()
    
    right_title_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.25), Inches(6), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    right_content = slide.shapes.add_textbox(Inches(6.8), Inches(1.9), Inches(5.8), Inches(5))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✓ {point}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    return slide

def add_comparison_table_slide(prs, title, headers, rows, highlight_col=0):
    """Create a comparison table slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    
    # Top accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = INFUSE_DEEP_ORANGE
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header
    
    table_width = Inches(12)
    table_height = Inches(5.5)
    left = Inches(0.67)
    top = Inches(1.3)
    
    table = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height).table
    
    # Set column widths
    col_width = table_width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = INFUSE_ORANGE if i == highlight_col else DARK_GRAY
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(14)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = value
            if col_idx == highlight_col:
                cell.fill.solid()
                cell.fill.fore_color.rgb = INFUSE_GOLD
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = DARK_GRAY
            para.alignment = PP_ALIGN.CENTER
    
    return slide

def create_infuse_pitch_deck():
    """Create the complete Infuse pitch deck"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # ========== SLIDE 1: Title Slide ==========
    add_title_slide(prs, "INFUSE", "Enterprise Healthcare & IoT Security Platform")
    
    # ========== SLIDE 2: Company Overview ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, DARK_GRAY, DARK_GRAY)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Who We Are"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = INFUSE_GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Mission statement
    mission_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.33), Inches(1.5))
    tf = mission_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Infuse is building enterprise-grade platforms that make healthcare accessible to every citizen and IoT ecosystems secure by default."
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Stats boxes
    stats = [("50M+", "Devices Secured"), ("500+", "Enterprise Clients"), ("99.9%", "Platform Uptime"), ("25+", "Countries")]
    for i, (value, label) in enumerate(stats):
        left = Inches(1 + i * 3)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(4), Inches(2.5), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = INFUSE_ORANGE if i % 2 == 0 else INFUSE_DEEP_ORANGE
        box.line.fill.background()
        
        value_box = slide.shapes.add_textbox(left, Inches(4.2), Inches(2.5), Inches(0.8))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        label_box = slide.shapes.add_textbox(left, Inches(5), Inches(2.5), Inches(0.6))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # ========== SLIDE 3: The Problem ==========
    add_content_slide(prs, "The Problem We're Solving", [
        "Healthcare systems are fragmented - patient data scattered across multiple providers",
        "Rural and underserved communities lack access to quality healthcare services",
        "IoT devices are multiplying exponentially but security lags far behind",
        "Telecom operators struggle to secure 5G networks and connected ecosystems",
        "Enterprise IoT deployments lack unified security visibility and threat response",
        "Existing solutions are either too expensive or too complex for mass adoption",
        "Data sovereignty regulations make global healthcare and IoT solutions challenging"
    ], INFUSE_DEEP_ORANGE)
    
    # ========== SLIDE 4: Our Solutions ==========
    add_two_column_slide(prs, "Two Platforms. Infinite Possibilities.",
        "HealthTrack Pro", [
            "AI-powered healthcare management platform",
            "Unified patient records across providers",
            "Real-time vitals monitoring & alerts",
            "Telemedicine with HD video consultations",
            "AI-assisted diagnosis suggestions",
            "ABDM/ABHA compliant (India)",
            "Multi-language support (8+ languages)",
            "Works offline in low-connectivity areas"
        ],
        "SecureSphere", [
            "Enterprise IoT & cybersecurity platform",
            "Real-time threat intelligence & detection",
            "Telecom-grade network security (5G/LTE)",
            "Automotive IoT security suite",
            "Smart city infrastructure protection",
            "CSP white-label & multi-tenant support",
            "Auto-discovery of network devices",
            "Compliance (SOC 2, HIPAA, ISO 27001)"
        ]
    )
    
    # ========== SLIDE 5: HealthTrack Pro USP ==========
    add_title_slide(prs, "HealthTrack Pro", "Transforming Healthcare Delivery")
    
    add_content_slide(prs, "HealthTrack Pro - Unique Selling Propositions", [
        "LAST-MILE HEALTHCARE: Designed to work in rural areas with low connectivity, offline-first architecture",
        "AI DIAGNOSTICS: ML-powered symptom analysis and diagnosis suggestions, reducing misdiagnosis by 40%",
        "UNIFIED PATIENT RECORDS: Single view of patient history across all healthcare providers",
        "ABHA/ABDM INTEGRATION: Fully compliant with India's Ayushman Bharat Digital Mission",
        "MULTI-LANGUAGE: Native support for Hindi, Tamil, Telugu, Bengali, and 8+ regional languages",
        "AFFORDABLE: 90% lower cost than competing enterprise healthcare solutions",
        "TELEMEDICINE BUILT-IN: HD video consultations with integrated prescription & follow-up management",
        "WEARABLE INTEGRATION: Real-time sync with smartwatches, BP monitors, glucose meters"
    ], INFUSE_ORANGE)
    
    # ========== SLIDE 6: HealthTrack Pro Future ==========
    add_content_slide(prs, "HealthTrack Pro - Future Roadmap", [
        "AI RADIOLOGY: Automated X-ray and CT scan analysis with 95%+ accuracy (Q2 2025)",
        "PREDICTIVE HEALTH: ML models predicting health risks 6-12 months in advance",
        "MENTAL HEALTH MODULE: AI-powered mental health screening and therapy tracking",
        "GENOMICS INTEGRATION: Personalized medicine based on genetic profiles",
        "PHARMACY NETWORK: Direct integration with pharmacy chains for medicine delivery",
        "INSURANCE AUTO-CLAIMS: Automated insurance claim processing and approval",
        "COMMUNITY HEALTH: Population health analytics for government health programs",
        "GLOBAL EXPANSION: FDA compliance for US market, EU MDR for European expansion"
    ], INFUSE_GOLD)
    
    # ========== SLIDE 7: HealthTrack Competition ==========
    add_comparison_table_slide(prs, "HealthTrack Pro vs Competition - Healthcare",
        ["Feature", "Infuse HealthTrack", "Practo", "1mg/Tata Health", "Apollo 24/7"],
        [
            ["AI Diagnostics", "✓ Advanced ML", "Basic", "✗ None", "Basic"],
            ["Offline Mode", "✓ Full Support", "✗ None", "✗ None", "✗ None"],
            ["Rural Connectivity", "✓ Optimized", "✗ Urban Only", "✗ Urban Only", "✗ Urban Only"],
            ["ABDM/ABHA", "✓ Full Compliance", "Partial", "Partial", "Partial"],
            ["Multi-Language", "✓ 8+ Languages", "3 Languages", "2 Languages", "3 Languages"],
            ["Wearable Sync", "✓ All Major Brands", "Limited", "✗ None", "Limited"],
            ["White-Label", "✓ Enterprise Ready", "✗ None", "✗ None", "✗ None"],
            ["Pricing", "90% Lower", "Premium", "Premium", "Premium"]
        ], 1)
    
    # ========== SLIDE 8: SecureSphere USP ==========
    add_title_slide(prs, "SecureSphere", "Securing the Connected World")
    
    add_content_slide(prs, "SecureSphere - Unique Selling Propositions", [
        "TELECOM-NATIVE: Purpose-built for CSPs with RAN sync, eSIM provisioning, and 5G security",
        "AI THREAT DETECTION: Real-time anomaly detection across millions of IoT endpoints",
        "CSP WHITE-LABEL: Full multi-tenant platform for telecom operators to resell",
        "AUTO-DISCOVERY: Automatic detection and profiling of all network devices",
        "NETFLOW ANALYTICS: Deep packet inspection and traffic analysis at scale",
        "AUTOMOTIVE SECURITY: Specialized modules for connected car ecosystems",
        "SMART CITY READY: Traffic cameras, smart meters, infrastructure protection",
        "UNIFIED DASHBOARD: Single pane of glass for all security operations"
    ], INFUSE_DEEP_ORANGE)
    
    # ========== SLIDE 9: SecureSphere Future ==========
    add_content_slide(prs, "SecureSphere - Future Roadmap", [
        "QUANTUM-SAFE ENCRYPTION: Post-quantum cryptography for future-proof security (Q3 2025)",
        "AI SOC AUTOMATION: Fully automated Security Operations Center with AI agents",
        "SATELLITE IoT: Security for satellite-connected IoT devices (LEO networks)",
        "DIGITAL TWIN SECURITY: Virtual replica security testing before production deployment",
        "ZERO-DAY PREDICTION: ML models predicting vulnerabilities before they're exploited",
        "BLOCKCHAIN AUDIT: Immutable audit trails for compliance and forensics",
        "EDGE AI: On-device threat detection for ultra-low latency response",
        "GLOBAL THREAT NETWORK: Shared threat intelligence across customer base"
    ], INFUSE_DARK_ORANGE)
    
    # ========== SLIDE 10: SecureSphere Competition ==========
    add_comparison_table_slide(prs, "SecureSphere vs Competition - IoT Security",
        ["Feature", "Infuse SecureSphere", "Palo Alto IoT", "Cisco IoT Threat", "Fortinet"],
        [
            ["CSP White-Label", "✓ Full Multi-Tenant", "✗ None", "✗ None", "Limited"],
            ["RAN/5G Security", "✓ Native Support", "Add-on", "Partial", "Add-on"],
            ["eSIM Provisioning", "✓ Built-in", "✗ None", "✗ None", "✗ None"],
            ["Auto-Discovery", "✓ AI-Powered", "Basic", "Basic", "Basic"],
            ["Automotive Suite", "✓ Specialized", "Generic", "✗ None", "Generic"],
            ["Smart City Module", "✓ Comprehensive", "Limited", "Limited", "Limited"],
            ["Deployment", "Cloud/On-Prem/Hybrid", "Cloud Only", "On-Prem Heavy", "Appliance"],
            ["Pricing Model", "Usage-Based", "Per-Device High", "Enterprise Only", "Per-Device High"]
        ], 1)
    
    # ========== SLIDE 11: Real Problems Solved ==========
    add_content_slide(prs, "Real-World Problems We Solve", [
        "HEALTHCARE ACCESS: 600M+ Indians lack access to quality healthcare - we bring it to their phones",
        "PATIENT DATA SILOS: Average patient visits 5+ providers - we unify their complete health history",
        "DELAYED DIAGNOSIS: Rural areas wait weeks for specialist opinions - we provide AI-assisted instant analysis",
        "IoT BLIND SPOTS: Enterprises have 30% unmanaged IoT devices - we discover and secure them all",
        "5G SECURITY GAP: Telecom operators lack native 5G security tools - we provide purpose-built solutions",
        "COMPLIANCE BURDEN: Healthcare & IoT face complex regulations - we automate compliance reporting",
        "LANGUAGE BARRIERS: Healthcare apps ignore regional languages - we support 8+ Indian languages",
        "COST PROHIBITIVE: Enterprise solutions price out SMBs - we offer 90% lower pricing"
    ], INFUSE_ORANGE)
    
    # ========== SLIDE 12: What Competitors Miss ==========
    add_two_column_slide(prs, "What Competitors Are Missing",
        "Healthcare Gap", [
            "No offline-first architecture for rural areas",
            "Limited regional language support",
            "No AI-powered diagnosis assistance",
            "Fragmented patient records",
            "No wearable device integration",
            "Urban-centric design philosophy",
            "Expensive enterprise pricing",
            "No government health program integration"
        ],
        "IoT Security Gap", [
            "No native telecom/CSP support",
            "Can't white-label for resale",
            "Manual device discovery only",
            "No automotive-specific security",
            "Limited smart city coverage",
            "No eSIM provisioning capability",
            "Rigid deployment options",
            "Per-device pricing kills scale"
        ]
    )
    
    # ========== SLIDE 13: Business Model ==========
    add_content_slide(prs, "Business Model", [
        "HEALTHTRACK PRO - SaaS subscription per healthcare provider + per patient premium features",
        "SECURESPHERE - Usage-based pricing (devices monitored, data processed, API calls)",
        "CSP WHITE-LABEL - Revenue share model with telecom operators (typically 70-30 split)",
        "ENTERPRISE LICENSES - Annual contracts for large hospital chains and corporations",
        "GOVERNMENT CONTRACTS - Per-citizen pricing for national health programs",
        "PROFESSIONAL SERVICES - Implementation, customization, and training",
        "DATA INSIGHTS - Anonymized, aggregated health and security analytics (opt-in)"
    ], INFUSE_DEEP_ORANGE)
    
    # ========== SLIDE 14: Traction ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, DARK_GRAY, DARK_GRAY)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Traction & Milestones"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = INFUSE_GOLD
    p.alignment = PP_ALIGN.CENTER
    
    milestones = [
        ("2019", "Company Founded"),
        ("2020", "HealthTrack MVP Launch"),
        ("2021", "First 100 Healthcare Partners"),
        ("2022", "SecureSphere Platform Launch"),
        ("2023", "500+ Enterprise Clients"),
        ("2024", "CSP White-Label Program"),
        ("2025", "Global Expansion Phase")
    ]
    
    for i, (year, milestone) in enumerate(milestones):
        top = Inches(1.8 + i * 0.75)
        
        # Year badge
        year_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), top, Inches(1.2), Inches(0.5))
        year_box.fill.solid()
        year_box.fill.fore_color.rgb = INFUSE_ORANGE
        year_box.line.fill.background()
        
        year_text = slide.shapes.add_textbox(Inches(1), top, Inches(1.2), Inches(0.5))
        tf = year_text.text_frame
        tf.paragraphs[0].text = year
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Milestone text
        milestone_text = slide.shapes.add_textbox(Inches(2.5), top, Inches(9), Inches(0.5))
        tf = milestone_text.text_frame
        tf.paragraphs[0].text = milestone
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = WHITE
    
    # ========== SLIDE 15: Team ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, DARK_GRAY, DARK_GRAY)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Leadership Team"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = INFUSE_GOLD
    p.alignment = PP_ALIGN.CENTER
    
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(0.5))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "74+ Years of Combined Experience"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    team = [
        ("Rohini Koul", "CEO", "20 years in Academia", "MSc Chemistry"),
        ("Chief Growth Officer", "Founder", "30 years in Technology", "BE, BIT | MBA, Arizona"),
        ("Chief Technology Officer", "CTO", "24 years in SaaS/PaaS", "BE, BIT | MBA, UPenn")
    ]
    
    for i, (name, role, exp, edu) in enumerate(team):
        left = Inches(0.5 + i * 4.2)
        
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(4), Inches(5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x2D, 0x33, 0x3F)
        card.line.fill.background()
        
        # Role badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(1), Inches(2), Inches(2), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = INFUSE_ORANGE
        badge.line.fill.background()
        
        badge_text = slide.shapes.add_textbox(left + Inches(1), Inches(2), Inches(2), Inches(0.4))
        tf = badge_text.text_frame
        tf.paragraphs[0].text = role
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Name
        name_box = slide.shapes.add_textbox(left, Inches(2.6), Inches(4), Inches(0.6))
        tf = name_box.text_frame
        tf.paragraphs[0].text = name
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Experience
        exp_box = slide.shapes.add_textbox(left, Inches(3.3), Inches(4), Inches(0.8))
        tf = exp_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = exp
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = INFUSE_GOLD
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Education
        edu_box = slide.shapes.add_textbox(left, Inches(4), Inches(4), Inches(0.8))
        tf = edu_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = edu
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ========== SLIDE 16: Contact ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, DARK_GRAY, DARK_GRAY)
    
    # Accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3), Inches(13.33), Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = INFUSE_ORANGE
    accent.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.33), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Let's Transform Together"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    contact_info = [
        "info@infuse.net.in",
        "+91-9599960663",
        "www.infuse.net.in",
        "India"
    ]
    
    for i, info in enumerate(contact_info):
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4 + i * 0.6), Inches(12.33), Inches(0.5))
        tf = info_box.text_frame
        p = tf.paragraphs[0]
        p.text = info
        p.font.size = Pt(20)
        p.font.color.rgb = INFUSE_GOLD if i == 0 else WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = "/app/frontend/public/Infuse_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"Pitch deck saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_infuse_pitch_deck()
