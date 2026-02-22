"""
Integration Architecture PPT Generator
Creates downloadable PowerPoint with block diagrams for:
1. SecureSphere CSP Integration
2. HealthTrack Pro Wearable/Health App Integration
"""
from fastapi import APIRouter
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import nsmap
import io

router = APIRouter(prefix="/api/integration-docs", tags=["Integration Documentation"])

# Colors
ORANGE = RGBColor(249, 115, 22)
DARK_ORANGE = RGBColor(234, 88, 12)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(30, 41, 59)
GRAY = RGBColor(55, 65, 81)
LIGHT_GRAY = RGBColor(243, 244, 246)
GREEN = RGBColor(16, 185, 129)
BLUE = RGBColor(59, 130, 246)
PURPLE = RGBColor(139, 92, 246)
RED = RGBColor(239, 68, 68)
AMBER = RGBColor(245, 158, 11)


def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_GRAY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
    tf = sub_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_section_slide(prs, title, color=ORANGE):
    """Add a section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_block(slide, left, top, width, height, text, fill_color, text_color=WHITE, font_size=12):
    """Add a block/box shape with text"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(5)
    
    return shape


def add_arrow_line(slide, start_x, start_y, end_x, end_y, color=GRAY):
    """Add a connector line"""
    connector = slide.shapes.add_connector(
        1,  # Straight connector
        Inches(start_x), Inches(start_y),
        Inches(end_x), Inches(end_y)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector


def create_csp_integration_slides(prs):
    """Create SecureSphere CSP Integration slides"""
    
    # Section Title
    add_section_slide(prs, "SecureSphere\nCSP Integration Architecture", BLUE)
    
    # Slide: CSP High-Level Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "CSP Integration: High-Level Architecture"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # CSP Layer (Left)
    add_block(slide, Inches(0.3), Inches(1), Inches(2), Inches(0.6), "CSP/Telecom Provider", PURPLE, WHITE, 11)
    add_block(slide, Inches(0.3), Inches(1.8), Inches(2), Inches(0.5), "Subscriber Base", GRAY, WHITE, 10)
    add_block(slide, Inches(0.3), Inches(2.4), Inches(2), Inches(0.5), "Network Infrastructure", GRAY, WHITE, 10)
    add_block(slide, Inches(0.3), Inches(3.0), Inches(2), Inches(0.5), "Billing System", GRAY, WHITE, 10)
    
    # Integration Layer (Middle)
    add_block(slide, Inches(3), Inches(1), Inches(2), Inches(0.6), "Integration Layer", ORANGE, WHITE, 11)
    add_block(slide, Inches(3), Inches(1.8), Inches(2), Inches(0.5), "REST API Gateway", AMBER, WHITE, 10)
    add_block(slide, Inches(3), Inches(2.4), Inches(2), Inches(0.5), "SDK (Android/iOS)", AMBER, WHITE, 10)
    add_block(slide, Inches(3), Inches(3.0), Inches(2), Inches(0.5), "Webhook Events", AMBER, WHITE, 10)
    add_block(slide, Inches(3), Inches(3.6), Inches(2), Inches(0.5), "White-Label Portal", AMBER, WHITE, 10)
    
    # SecureSphere Core (Right)
    add_block(slide, Inches(5.7), Inches(1), Inches(2.2), Inches(0.6), "SecureSphere Core", BLUE, WHITE, 11)
    add_block(slide, Inches(5.7), Inches(1.8), Inches(2.2), Inches(0.5), "URL Scanner Engine", DARK_GRAY, WHITE, 10)
    add_block(slide, Inches(5.7), Inches(2.4), Inches(2.2), Inches(0.5), "SMS Fraud Detection", DARK_GRAY, WHITE, 10)
    add_block(slide, Inches(5.7), Inches(3.0), Inches(2.2), Inches(0.5), "Threat Intelligence", DARK_GRAY, WHITE, 10)
    add_block(slide, Inches(5.7), Inches(3.6), Inches(2.2), Inches(0.5), "AI/ML Models", DARK_GRAY, WHITE, 10)
    
    # End Users (Far Right)
    add_block(slide, Inches(8.2), Inches(1), Inches(1.5), Inches(0.6), "End Users", GREEN, WHITE, 11)
    add_block(slide, Inches(8.2), Inches(1.8), Inches(1.5), Inches(0.5), "Mobile App", RGBColor(34, 197, 94), WHITE, 10)
    add_block(slide, Inches(8.2), Inches(2.4), Inches(1.5), Inches(0.5), "SMS Alerts", RGBColor(34, 197, 94), WHITE, 10)
    add_block(slide, Inches(8.2), Inches(3.0), Inches(1.5), Inches(0.5), "Web Portal", RGBColor(34, 197, 94), WHITE, 10)
    
    # Data Flow Description
    desc_box = slide.shapes.add_textbox(Inches(0.3), Inches(4.4), Inches(9.4), Inches(2.5))
    tf = desc_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Data Flow:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    flows = [
        "1. CSP integrates via REST API or SDK → SecureSphere processes requests",
        "2. Real-time threat detection for URLs, SMS, calls → Results returned via API/Webhook",
        "3. CSP white-labels the service → End users access via CSP-branded app/portal",
        "4. Usage metrics sent to CSP billing → Seamless subscription management"
    ]
    
    for flow in flows:
        p = tf.add_paragraph()
        p.text = flow
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.space_before = Pt(4)
    
    # Slide: CSP API Integration Flow
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = WHITE
    
    title_box = slide2.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "CSP API Integration: Request Flow"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Flow boxes
    add_block(slide2, Inches(0.5), Inches(1.2), Inches(1.8), Inches(0.8), "CSP Backend\nSystem", PURPLE, WHITE, 11)
    add_block(slide2, Inches(2.8), Inches(1.2), Inches(1.8), Inches(0.8), "API Gateway\n(Auth + Rate Limit)", ORANGE, WHITE, 10)
    add_block(slide2, Inches(5.1), Inches(1.2), Inches(1.8), Inches(0.8), "SecureSphere\nEngine", BLUE, WHITE, 11)
    add_block(slide2, Inches(7.4), Inches(1.2), Inches(1.8), Inches(0.8), "Response\n+ Webhook", GREEN, WHITE, 11)
    
    # API Endpoints
    endpoints_box = slide2.shapes.add_textbox(Inches(0.3), Inches(2.3), Inches(9.4), Inches(4))
    tf = endpoints_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Available API Endpoints for CSPs:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    endpoints = [
        ("POST /api/securesphere/url/scan", "Scan URL for threats (phishing, malware, scam)"),
        ("POST /api/securesphere/url/scan/bulk", "Bulk URL scanning (up to 100 URLs)"),
        ("POST /api/securesphere/sms/analyze", "Analyze SMS for fraud patterns"),
        ("POST /api/securesphere/sms/analyze/bulk", "Bulk SMS analysis"),
        ("GET /api/securesphere/threat-score/{entity}", "Get threat score for entity"),
        ("GET /api/securesphere/device/registry", "Device inventory management"),
        ("POST /api/securesphere/webhook/register", "Register webhook for real-time alerts"),
        ("GET /api/securesphere/analytics/usage", "Usage analytics and reporting"),
    ]
    
    for endpoint, desc in endpoints:
        p = tf.add_paragraph()
        p.text = f"• {endpoint}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.space_before = Pt(6)
        
        p = tf.add_paragraph()
        p.text = f"   {desc}"
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY
    
    # Slide: CSP Use Cases
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = WHITE
    
    title_box = slide3.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "SecureSphere: CSP Use Cases"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    use_cases = [
        ("SMS Fraud Protection", "Filter scam/phishing SMS before delivery to subscribers", BLUE),
        ("Safe Browsing Service", "Real-time URL checking in CSP browser/app", GREEN),
        ("Call Spam Detection", "Identify and block spam/scam calls", PURPLE),
        ("IoT Device Security", "Secure smart home devices on CSP network", ORANGE),
        ("5G Network Security", "Protect 5G infrastructure and edge devices", RED),
        ("Enterprise Fleet Protection", "Secure corporate device fleets", AMBER),
    ]
    
    y_pos = 1.0
    for i, (title, desc, color) in enumerate(use_cases):
        col = i % 2
        row = i // 2
        x = 0.5 + (col * 4.8)
        y = 1.0 + (row * 1.8)
        
        add_block(slide3, Inches(x), Inches(y), Inches(4.5), Inches(0.6), title, color, WHITE, 14)
        
        desc_box = slide3.shapes.add_textbox(Inches(x), Inches(y + 0.65), Inches(4.5), Inches(0.9))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY


def create_securesphere_all_usecases(prs):
    """Create SecureSphere all use cases slides"""
    
    add_section_slide(prs, "SecureSphere\nComplete Use Case Overview", DARK_GRAY)
    
    # Slide: All SecureSphere Verticals
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "SecureSphere: Multi-Vertical Security Platform"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Core Platform (Center)
    add_block(slide, Inches(3.5), Inches(2.8), Inches(3), Inches(1), "SecureSphere\nAI Security Core", ORANGE, WHITE, 14)
    
    # Verticals around it
    verticals = [
        (1.5, 0.9, "Consumer\nMobile Security", BLUE),
        (5.5, 0.9, "Enterprise\nFleet Management", GREEN),
        (0.3, 2.8, "Telecom/CSP\nIntegration", PURPLE),
        (7.2, 2.8, "Automotive\nConnected Vehicles", RED),
        (1.5, 4.8, "IoT/Smart Home\nProtection", AMBER),
        (5.5, 4.8, "5G Network\nSecurity", RGBColor(6, 182, 212)),
    ]
    
    for x, y, text, color in verticals:
        add_block(slide, Inches(x), Inches(y), Inches(2.5), Inches(1), text, color, WHITE, 11)
    
    # Slide: Automotive Use Case Detail
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = WHITE
    
    title_box = slide2.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "Automotive Security: Connected Vehicle Protection"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Vehicle
    add_block(slide2, Inches(0.5), Inches(1.2), Inches(2.5), Inches(1.2), "Connected\nVehicle", RED, WHITE, 14)
    
    # Vehicle components
    components = ["ECU Monitoring", "CAN Bus Security", "OTA Updates", "V2X Communication"]
    for i, comp in enumerate(components):
        add_block(slide2, Inches(0.5), Inches(2.6 + i*0.55), Inches(2.5), Inches(0.45), comp, GRAY, WHITE, 10)
    
    # SecureSphere
    add_block(slide2, Inches(3.8), Inches(2), Inches(2.4), Inches(1.5), "SecureSphere\nAutomotive\nSecurity", ORANGE, WHITE, 12)
    
    # OEM/Fleet
    add_block(slide2, Inches(7), Inches(1.2), Inches(2.5), Inches(1.2), "OEM / Fleet\nManager", BLUE, WHITE, 14)
    
    features = ["Real-time Alerts", "Security Dashboard", "Compliance Reports", "Threat Analytics"]
    for i, feat in enumerate(features):
        add_block(slide2, Inches(7), Inches(2.6 + i*0.55), Inches(2.5), Inches(0.45), feat, GRAY, WHITE, 10)
    
    # Slide: IoT/5G Security
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = WHITE
    
    title_box = slide3.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "5G-IoT Security: Network-Wide Protection"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # IoT Devices
    add_block(slide3, Inches(0.3), Inches(1), Inches(2), Inches(0.5), "Smart Home", GREEN, WHITE, 11)
    add_block(slide3, Inches(0.3), Inches(1.6), Inches(2), Inches(0.5), "Industrial IoT", GREEN, WHITE, 11)
    add_block(slide3, Inches(0.3), Inches(2.2), Inches(2), Inches(0.5), "Wearables", GREEN, WHITE, 11)
    add_block(slide3, Inches(0.3), Inches(2.8), Inches(2), Inches(0.5), "Smart Cities", GREEN, WHITE, 11)
    
    # 5G Network
    add_block(slide3, Inches(2.8), Inches(1.5), Inches(2), Inches(1.5), "5G\nNetwork\nInfrastructure", PURPLE, WHITE, 12)
    
    # SecureSphere 5G
    add_block(slide3, Inches(5.3), Inches(1), Inches(2.3), Inches(0.6), "SecureSphere 5G-IoT", ORANGE, WHITE, 11)
    add_block(slide3, Inches(5.3), Inches(1.7), Inches(2.3), Inches(0.45), "Network Slicing Security", AMBER, WHITE, 9)
    add_block(slide3, Inches(5.3), Inches(2.2), Inches(2.3), Inches(0.45), "Edge Computing Protection", AMBER, WHITE, 9)
    add_block(slide3, Inches(5.3), Inches(2.7), Inches(2.3), Inches(0.45), "Device Authentication", AMBER, WHITE, 9)
    add_block(slide3, Inches(5.3), Inches(3.2), Inches(2.3), Inches(0.45), "Threat Intelligence", AMBER, WHITE, 9)
    
    # Compliance
    add_block(slide3, Inches(8), Inches(1), Inches(1.7), Inches(0.5), "NIST CSF", BLUE, WHITE, 10)
    add_block(slide3, Inches(8), Inches(1.6), Inches(1.7), Inches(0.5), "ISO 27001", BLUE, WHITE, 10)
    add_block(slide3, Inches(8), Inches(2.2), Inches(1.7), Inches(0.5), "IEC 62443", BLUE, WHITE, 10)
    add_block(slide3, Inches(8), Inches(2.8), Inches(1.7), Inches(0.5), "ETSI EN 303", BLUE, WHITE, 10)
    
    # Description
    desc_box = slide3.shapes.add_textbox(Inches(0.3), Inches(4), Inches(9.4), Inches(2.5))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Capabilities: Real-time breach detection • Quantum-safe encryption readiness • Deep network visibility • AI-powered threat analysis • Multi-protocol support (NB-IoT, LTE-M, 5G-NR)"
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY


def create_healthtrack_integration_slides(prs):
    """Create HealthTrack Pro wearable integration slides"""
    
    add_section_slide(prs, "HealthTrack Pro\nWearable & Health App Integration", GREEN)
    
    # Slide: Wearable Integration Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "HealthTrack Pro: Universal Health Data Integration"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Data Sources (Left)
    add_block(slide, Inches(0.3), Inches(0.9), Inches(2.2), Inches(0.5), "Data Sources", PURPLE, WHITE, 12)
    
    sources = [
        "Apple HealthKit", "Google Fit", "Samsung Health", 
        "Fitbit API", "Garmin Connect", "Withings API",
        "FHIR/HL7 Systems", "Any REST API"
    ]
    for i, src in enumerate(sources):
        add_block(slide, Inches(0.3), Inches(1.5 + i*0.45), Inches(2.2), Inches(0.4), src, GRAY, WHITE, 9)
    
    # Integration Hub (Center)
    add_block(slide, Inches(3.2), Inches(1.5), Inches(3.2), Inches(0.6), "HealthTrack Integration Hub", ORANGE, WHITE, 12)
    
    hub_features = [
        "Data Normalization Engine",
        "Protocol Adapters (FHIR, HL7, REST)",
        "Real-time Sync Manager",
        "Privacy & Consent Handler",
        "AI Health Analytics"
    ]
    for i, feat in enumerate(hub_features):
        add_block(slide, Inches(3.2), Inches(2.2 + i*0.5), Inches(3.2), Inches(0.45), feat, AMBER, WHITE, 9)
    
    # Output (Right)
    add_block(slide, Inches(7), Inches(0.9), Inches(2.7), Inches(0.5), "HealthTrack Pro", GREEN, WHITE, 12)
    
    outputs = [
        "Unified Patient Dashboard",
        "AI Health Insights",
        "Predictive Analytics",
        "Doctor Notifications",
        "EHR Integration"
    ]
    for i, out in enumerate(outputs):
        add_block(slide, Inches(7), Inches(1.5 + i*0.5), Inches(2.7), Inches(0.45), out, RGBColor(34, 197, 94), WHITE, 9)
    
    # Slide: OEM Wearable Integration
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = WHITE
    
    title_box = slide2.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "OEM Wearable Integration: Connect Any Device"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # OEM Device
    add_block(slide2, Inches(0.5), Inches(1.2), Inches(2.5), Inches(1), "OEM Wearable\nDevice", PURPLE, WHITE, 14)
    
    # Data types
    data_types = ["Heart Rate", "SpO2", "Steps/Activity", "Sleep Data", "Blood Pressure", "ECG/PPG"]
    for i, dt in enumerate(data_types):
        col = i % 2
        row = i // 2
        add_block(slide2, Inches(0.5 + col*1.3), Inches(2.4 + row*0.5), Inches(1.2), Inches(0.45), dt, GRAY, WHITE, 9)
    
    # Integration Options
    add_block(slide2, Inches(3.5), Inches(1.2), Inches(2.8), Inches(0.6), "Integration Options", ORANGE, WHITE, 12)
    
    options = [
        ("SDK Integration", "Native iOS/Android SDK"),
        ("REST API", "Direct HTTP API calls"),
        ("BLE Gateway", "Bluetooth data bridge"),
        ("Cloud-to-Cloud", "Backend API sync"),
    ]
    y = 1.9
    for title, desc in options:
        add_block(slide2, Inches(3.5), Inches(y), Inches(2.8), Inches(0.4), f"{title}: {desc}", AMBER, WHITE, 9)
        y += 0.45
    
    # HealthTrack Features
    add_block(slide2, Inches(7), Inches(1.2), Inches(2.7), Inches(0.6), "HealthTrack Features", GREEN, WHITE, 12)
    
    features = [
        "Real-time Vitals Monitoring",
        "AI Anomaly Detection",
        "Trend Analysis",
        "Doctor Alerts",
        "Patient Reports",
        "Prescription Correlation"
    ]
    for i, feat in enumerate(features):
        add_block(slide2, Inches(7), Inches(1.9 + i*0.45), Inches(2.7), Inches(0.4), feat, RGBColor(34, 197, 94), WHITE, 9)
    
    # Slide: Health App Integration
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = WHITE
    
    title_box = slide3.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "Phone Health App Integration: Data Flow"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Phone Apps
    add_block(slide3, Inches(0.5), Inches(1), Inches(2.3), Inches(0.6), "Phone Health Apps", BLUE, WHITE, 12)
    
    apps = [
        ("Apple Health", "HealthKit API"),
        ("Google Fit", "Fitness REST API"),
        ("Samsung Health", "Health SDK"),
        ("Huawei Health", "Health Kit"),
    ]
    y = 1.7
    for app, api in apps:
        add_block(slide3, Inches(0.5), Inches(y), Inches(2.3), Inches(0.5), f"{app}\n{api}", GRAY, WHITE, 9)
        y += 0.6
    
    # HealthTrack Mobile SDK
    add_block(slide3, Inches(3.5), Inches(1.5), Inches(3), Inches(1.5), "HealthTrack\nMobile SDK\n\n• Auto-sync\n• Offline Support\n• Encryption", ORANGE, WHITE, 10)
    
    # Cloud Processing
    add_block(slide3, Inches(3.5), Inches(3.3), Inches(3), Inches(1.2), "Cloud Processing\n\n• Data Validation\n• AI Analysis\n• Storage", AMBER, WHITE, 10)
    
    # Doctor Dashboard
    add_block(slide3, Inches(7.2), Inches(1), Inches(2.5), Inches(0.6), "Doctor Dashboard", GREEN, WHITE, 12)
    
    doctor_features = [
        "Patient Vitals Overview",
        "AI Health Insights",
        "Trend Alerts",
        "Prescription Impact",
        "Video Consultation"
    ]
    for i, feat in enumerate(doctor_features):
        add_block(slide3, Inches(7.2), Inches(1.7 + i*0.5), Inches(2.5), Inches(0.45), feat, RGBColor(34, 197, 94), WHITE, 9)
    
    # Slide: API Endpoints for Health Data
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    slide4.background.fill.solid()
    slide4.background.fill.fore_color.rgb = WHITE
    
    title_box = slide4.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "HealthTrack Pro: Integration API Endpoints"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    endpoints_box = slide4.shapes.add_textbox(Inches(0.3), Inches(0.9), Inches(9.4), Inches(5.5))
    tf = endpoints_box.text_frame
    tf.word_wrap = True
    
    endpoints = [
        ("POST /api/healthtrack/vitals/sync", "Sync vitals from wearable/app"),
        ("POST /api/healthtrack/vitals/bulk", "Bulk vitals upload"),
        ("GET /api/healthtrack/patients/{id}/vitals", "Get patient vitals history"),
        ("POST /api/healthtrack/wearable/register", "Register new wearable device"),
        ("GET /api/healthtrack/wearable/supported", "List supported devices"),
        ("POST /api/healthtrack/healthkit/sync", "Sync from Apple HealthKit"),
        ("POST /api/healthtrack/googlefit/sync", "Sync from Google Fit"),
        ("GET /api/healthtrack/ai/health-insights", "Get AI health analysis"),
        ("POST /api/healthtrack/webhook/vitals", "Webhook for real-time vitals"),
        ("GET /api/healthtrack/analytics/trends", "Patient health trends"),
    ]
    
    for endpoint, desc in endpoints:
        p = tf.add_paragraph()
        p.text = f"• {endpoint}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = GREEN
        p.space_before = Pt(6)
        
        p = tf.add_paragraph()
        p.text = f"   {desc}"
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY


def create_summary_slide(prs):
    """Create summary slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_GRAY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = "Ready to Integrate?"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(3))
    tf = contact_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "SecureSphere & HealthTrack Pro"
    p.font.size = Pt(28)
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\nwww.infuse.net.in"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\nAPI Documentation • SDK Downloads • Developer Portal"
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER


@router.get("/architecture-ppt")
async def generate_integration_architecture_ppt():
    """Generate and download Integration Architecture PowerPoint"""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    add_title_slide(prs, "Infuse-AI Integration Architecture", 
                   "SecureSphere & HealthTrack Pro • Integration Guide")
    
    # SecureSphere CSP Integration
    create_csp_integration_slides(prs)
    
    # SecureSphere All Use Cases
    create_securesphere_all_usecases(prs)
    
    # HealthTrack Pro Wearable Integration
    create_healthtrack_integration_slides(prs)
    
    # Summary
    create_summary_slide(prs)
    
    # Save to bytes
    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    
    return StreamingResponse(
        ppt_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": "attachment; filename=Infuse-AI-Integration-Architecture.pptx"
        }
    )
