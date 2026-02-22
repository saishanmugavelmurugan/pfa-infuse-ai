"""
HealthTrack Pro Pitch Deck Generator
Comprehensive pitch deck focused on HealthTrack Pro with detailed features for testing
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import requests
from io import BytesIO

# Infuse Brand Colors
INFUSE_GOLD = RGBColor(0xFF, 0xDA, 0x7B)      # #FFDA7B
INFUSE_ORANGE = RGBColor(0xFF, 0x9A, 0x3B)    # #FF9A3B
INFUSE_DEEP_ORANGE = RGBColor(0xE5, 0x5A, 0x00)  # #E55A00
INFUSE_DARK_ORANGE = RGBColor(0xC6, 0x47, 0x00)  # #C64700
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x1F, 0x29, 0x37)
LIGHT_GRAY = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)

def add_logo(slide, logo_path="/app/backend/utils/infuse_logo.jpg"):
    """Add Infuse logo to top-right corner"""
    try:
        # Add actual logo image to top-right corner
        logo = slide.shapes.add_picture(logo_path, Inches(11.3), Inches(0.15), height=Inches(0.6))
    except Exception as e:
        # Fallback to text if image fails
        logo_box = slide.shapes.add_textbox(Inches(11.3), Inches(0.2), Inches(1.8), Inches(0.4))
        tf = logo_box.text_frame
        p = tf.paragraphs[0]
        p.text = "INFUSE"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = INFUSE_ORANGE

def add_dark_background(slide):
    """Add dark background"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_GRAY
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_white_background(slide):
    """Add white background"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_accent_bar(slide, color=INFUSE_ORANGE, top=0):
    """Add accent bar"""
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(top), Inches(13.33), Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()

def create_title_slide(prs, title, subtitle=""):
    """Create a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_logo(slide)
    add_accent_bar(slide, INFUSE_ORANGE, 3.5)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.33), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = INFUSE_GOLD
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_content_slide(prs, title, points, accent_color=INFUSE_ORANGE):
    """Create a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide)
    add_logo(slide)
    add_accent_bar(slide, accent_color, 0)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.33), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)
    
    return slide

def create_comparison_slide(prs, title, headers, rows):
    """Create a comparison table slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide)
    add_logo(slide)
    add_accent_bar(slide, INFUSE_DEEP_ORANGE, 0)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    num_cols = len(headers)
    num_rows = len(rows) + 1
    
    table_width = Inches(12.5)
    left = Inches(0.42)
    top = Inches(1)
    
    # Calculate row height based on number of rows
    max_rows = 12
    row_height = min(Inches(0.5), Inches(6) / num_rows)
    table_height = row_height * num_rows
    
    table = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height).table
    
    # Set column widths
    col_widths = [Inches(2.5)] + [Inches((12.5 - 2.5) / (num_cols - 1))] * (num_cols - 1)
    for i, width in enumerate(col_widths):
        table.columns[i].width = int(width)
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = INFUSE_ORANGE if i == 1 else DARK_GRAY
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(11)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            if col_idx == 1:  # HealthTrack column
                cell.fill.solid()
                cell.fill.fore_color.rgb = INFUSE_GOLD
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(10)
            para.font.color.rgb = DARK_GRAY
            para.alignment = PP_ALIGN.CENTER
    
    return slide

def create_feature_workflow_slide(prs, feature_name, steps, test_cases):
    """Create a feature workflow slide for testers"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide)
    add_logo(slide)
    add_accent_bar(slide, INFUSE_ORANGE, 0)
    
    # Feature name
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Feature: {feature_name}"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Workflow section
    workflow_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.9), Inches(6), Inches(0.35))
    workflow_header.fill.solid()
    workflow_header.fill.fore_color.rgb = INFUSE_ORANGE
    workflow_header.line.fill.background()
    
    wf_title = slide.shapes.add_textbox(Inches(0.3), Inches(0.9), Inches(6), Inches(0.35))
    tf = wf_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Workflow Steps"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    steps_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.35), Inches(6), Inches(5.5))
    tf = steps_box.text_frame
    tf.word_wrap = True
    for i, step in enumerate(steps):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{i+1}. {step}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    
    # Test cases section
    test_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(0.9), Inches(6), Inches(0.35))
    test_header.fill.solid()
    test_header.fill.fore_color.rgb = INFUSE_DEEP_ORANGE
    test_header.line.fill.background()
    
    tc_title = slide.shapes.add_textbox(Inches(6.7), Inches(0.9), Inches(6), Inches(0.35))
    tf = tc_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Test Cases"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    tests_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.35), Inches(6), Inches(5.5))
    tf = tests_box.text_frame
    tf.word_wrap = True
    for i, test in enumerate(test_cases):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"TC{i+1}: {test}"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(5)
    
    return slide

def create_roadmap_slide(prs, title, years_data):
    """Create a roadmap timeline slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_logo(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = INFUSE_GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Timeline
    colors = [INFUSE_GOLD, INFUSE_ORANGE, INFUSE_DEEP_ORANGE]
    for i, (year, items) in enumerate(years_data):
        left = Inches(0.5 + i * 4.2)
        
        # Year header
        year_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1), Inches(4), Inches(0.5))
        year_box.fill.solid()
        year_box.fill.fore_color.rgb = colors[i]
        year_box.line.fill.background()
        
        year_text = slide.shapes.add_textbox(left, Inches(1), Inches(4), Inches(0.5))
        tf = year_text.text_frame
        p = tf.paragraphs[0]
        p.text = year
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE if i > 0 else DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
        
        # Items
        items_box = slide.shapes.add_textbox(left, Inches(1.6), Inches(4), Inches(5.5))
        tf = items_box.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.space_after = Pt(6)
    
    return slide

def create_healthtrack_pitch_deck():
    """Create the HealthTrack Pro focused pitch deck"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # ========== SLIDE 1: Title ==========
    create_title_slide(prs, "HealthTrack Pro", "AI-Powered Lab Analysis | Wearable Integration | Ayurvedic Lifestyle Wellness")
    
    # ========== SLIDE 2: Vision ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_logo(slide)
    
    vision_label = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.5))
    tf = vision_label.text_frame
    p = tf.paragraphs[0]
    p.text = "OUR VISION"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = INFUSE_GOLD
    p.alignment = PP_ALIGN.CENTER
    
    vision_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.73), Inches(2.5))
    tf = vision_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "To revolutionize personal wellness by intelligently analyzing lab reports and wearable health data, providing AI-powered insights combined with personalized Ayurvedic lifestyle recommendations - taking India's 5000-year-old wellness wisdom to the entire world."
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Key pillars
    pillars = [
        ("Lab Report Analysis", "Pull & analyze reports from any lab"),
        ("Wearable Integration", "Apple Health, Samsung Health & more"),
        ("Ayurvedic Lifestyle", "Yoga, Diet, Daily Routine suggestions")
    ]
    
    for i, (pillar_title, pillar_desc) in enumerate(pillars):
        left = Inches(0.8 + i * 4.2)
        
        pillar_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(4.3), Inches(3.8), Inches(1.2))
        pillar_box.fill.solid()
        pillar_box.fill.fore_color.rgb = RGBColor(0x2D, 0x33, 0x3F)
        pillar_box.line.color.rgb = INFUSE_ORANGE
        
        p_title = slide.shapes.add_textbox(left, Inches(4.4), Inches(3.8), Inches(0.5))
        tf = p_title.text_frame
        p = tf.paragraphs[0]
        p.text = pillar_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = INFUSE_ORANGE
        p.alignment = PP_ALIGN.CENTER
        
        p_desc = slide.shapes.add_textbox(left, Inches(4.9), Inches(3.8), Inches(0.5))
        tf = p_desc.text_frame
        p = tf.paragraphs[0]
        p.text = pillar_desc
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    tagline = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.33), Inches(1))
    tf = tagline.text_frame
    p = tf.paragraphs[0]
    p.text = '"Taking Ayurveda to the World - Wellness, Not Just Healthcare"'
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = INFUSE_GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # ========== SLIDE 3: What We Do ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_logo(slide)
    
    mission_label = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.5))
    tf = mission_label.text_frame
    p = tf.paragraphs[0]
    p.text = "WHAT WE DO - NOT TEST BOOKING, BUT INTELLIGENT ANALYSIS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = INFUSE_GOLD
    p.alignment = PP_ALIGN.CENTER
    
    missions = [
        ("Pull Lab Reports", "We pull your lab reports directly from diagnostic labs - no manual uploads needed. Our AI extracts every parameter automatically."),
        ("Wearable Data Sync", "Connect Apple Health, Samsung Health, Google Fit - we analyze sleep, heart rate, HRV, SpO2, steps, stress levels & more."),
        ("AI-Powered Analysis", "Our AI correlates lab parameters with wearable data to generate comprehensive health insights and trends."),
        ("Ayurvedic Lifestyle Advice", "Based on your data, we suggest personalized Yoga asanas, Ayurvedic diet plans, and daily routines (Dinacharya)."),
        ("Clear Disclaimer", "All suggestions are lifestyle recommendations, NOT medical advice. We always recommend consulting qualified doctors."),
        ("Doctor Platform", "Verified Ayurvedic doctors (BAMS/MD) on platform for personalized consultations - taking Ayurveda global!")
    ]
    
    for i, (title, desc) in enumerate(missions):
        row = i // 2
        col = i % 2
        left = Inches(0.5 + col * 6.4)
        top = Inches(1.2 + row * 2)
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6), Inches(1.7))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x2D, 0x33, 0x3F)
        box.line.fill.background()
        
        title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), Inches(5.6), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = INFUSE_ORANGE
        
        desc_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.55), Inches(5.6), Inches(1.1))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
    
    # ========== SLIDE 4: Objectives ==========
    create_content_slide(prs, "Strategic Objectives 2025-2027", [
        "MARKET PENETRATION: Onboard 10,000+ healthcare providers across Tier 2/3 cities in Year 1",
        "PATIENT REACH: Enable healthcare access for 10 million patients by end of Year 2",
        "AI ACCURACY: Achieve 95%+ accuracy in AI-assisted diagnosis across 50+ common conditions",
        "LANGUAGE COVERAGE: Support all 22 scheduled Indian languages by Year 3",
        "GOVERNMENT PARTNERSHIPS: Integrate with 5+ state health programs for public healthcare delivery",
        "INTERNATIONAL EXPANSION: Launch in 3 Southeast Asian markets (Philippines, Indonesia, Vietnam) by Year 3",
        "REVENUE TARGET: Achieve ₹100 Cr ARR by Year 3 with 40% gross margins",
        "COMPLIANCE: Maintain 100% ABDM compliance and achieve ISO 27001 + SOC 2 certifications"
    ], INFUSE_ORANGE)
    
    # ========== SLIDE 5: Why Now? ==========
    create_content_slide(prs, "Why We Are Entering the Market NOW", [
        "DIGITAL INDIA MOMENTUM: 800M+ Indians now have smartphones; digital health adoption at all-time high",
        "ABDM ECOSYSTEM: Government's Ayushman Bharat Digital Mission creating unified health infrastructure",
        "POST-COVID AWARENESS: Pandemic accelerated telemedicine acceptance across all demographics",
        "RURAL CONNECTIVITY: 4G coverage now reaches 95%+ of India; 5G rolling out in urban areas",
        "HEALTHCARE CRISIS: India has only 1 doctor per 1,445 people (WHO recommends 1:1000) - technology must bridge this gap",
        "COMPETITOR GAP: Existing players (Practo, 1mg, Apollo) focus on urban metros; rural market underserved",
        "AI MATURITY: Large Language Models now accurate enough for reliable medical assistance",
        "FUNDING ENVIRONMENT: Digital health is top investment priority for VCs in India",
        "WEARABLES BOOM: Smartwatch penetration growing 50% YoY; enabling continuous health monitoring"
    ], INFUSE_DEEP_ORANGE)
    
    # ========== SLIDE 6: Market Comparison 1 ==========
    create_comparison_slide(prs, "HealthTrack Pro vs Competition - Feature Comparison",
        ["Feature", "HealthTrack Pro", "Practo", "1mg/Tata Health", "Apollo 24/7", "mFine", "Lybrate"],
        [
            ["AI Lab Analysis", "✓ Advanced ML", "✗ None", "✗ None", "Basic", "✗ None", "✗ None"],
            ["Lab Report Pull", "✓ Auto-Pull", "✗ Manual", "✗ Manual", "✗ Manual", "✗ Manual", "✗ Manual"],
            ["Ayurvedic Focus", "✓ Core Feature", "✗ None", "✗ None", "✗ None", "✗ None", "✗ None"],
            ["Wearable Sync", "✓ All Platforms", "Limited", "✗ None", "Limited", "✗ None", "✗ None"],
            ["Health Correlation", "✓ Lab+Wearable", "✗ None", "✗ None", "✗ None", "✗ None", "✗ None"],
            ["Lifestyle Suggestions", "✓ Yoga+Diet+Routine", "✗ None", "✗ None", "✗ None", "✗ None", "✗ None"],
            ["Languages", "8+", "3", "2", "3", "2", "2"],
            ["Ayurvedic Doctors", "✓ Verified", "✗ None", "✗ None", "✗ None", "✗ None", "✗ None"],
            ["Global Platform", "✓ Multi-currency", "India Only", "India Only", "India Only", "India Only", "India Only"],
            ["Offline Mode", "✓ Full Support", "✗ None", "✗ None", "✗ None", "✗ None", "✗ None"]
        ])
    
    # ========== SLIDE 7: Market Comparison 2 ==========
    create_comparison_slide(prs, "HealthTrack Pro vs Competition - Business & Technical",
        ["Aspect", "HealthTrack Pro", "Practo", "1mg", "Apollo 24/7", "mFine"],
        [
            ["Target Market", "Pan-India + Rural", "Urban Metros", "Urban + Semi", "Apollo Network", "Metro Cities"],
            ["Pricing Model", "Freemium + B2B", "Premium", "Premium", "Premium", "Premium"],
            ["Cost to Provider", "90% Lower", "High", "High", "Tied to Apollo", "High"],
            ["Data Ownership", "Provider Owns", "Platform Owns", "Platform Owns", "Apollo Owns", "Platform Owns"],
            ["API Access", "✓ Full APIs", "Limited", "✗ None", "✗ None", "✗ None"],
            ["Custom Branding", "✓ White-Label", "✗ No", "✗ No", "✗ No", "✗ No"],
            ["Govt Integration", "✓ Ready", "Partial", "Partial", "✗ No", "✗ No"],
            ["Patient Records", "Unified View", "Fragmented", "Fragmented", "Apollo Only", "Fragmented"],
            ["Deployment", "Cloud/On-Prem", "Cloud Only", "Cloud Only", "Cloud Only", "Cloud Only"]
        ])
    
    # ========== SLIDE 8: Differentiators ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide)
    add_logo(slide)
    add_accent_bar(slide, INFUSE_ORANGE, 0)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "How We Create Lasting Differentiation"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    differentiators = [
        ("OFFLINE-FIRST ARCHITECTURE", "Built from ground up for low/no connectivity. Full functionality without internet. Auto-sync when connection restored.", INFUSE_ORANGE),
        ("AI THAT LEARNS LOCALLY", "Our ML models trained on Indian health data patterns, understanding regional disease prevalence and symptoms.", INFUSE_DEEP_ORANGE),
        ("LANGUAGE-NATIVE DESIGN", "Not just translation - culturally adapted interfaces for each language with local health terminology.", INFUSE_DARK_ORANGE),
        ("PROVIDER ECONOMICS", "90% lower cost + revenue share model. Providers earn from the platform, not just pay for it.", INFUSE_ORANGE),
        ("DATA SOVEREIGNTY", "Patient data stays with providers. No platform lock-in. Full portability guaranteed.", INFUSE_DEEP_ORANGE),
        ("ECOSYSTEM INTEGRATION", "Only platform with full ABDM + state health program + insurance + pharmacy network integration.", INFUSE_DARK_ORANGE)
    ]
    
    for i, (title, desc, color) in enumerate(differentiators):
        row = i // 2
        col = i % 2
        left = Inches(0.4 + col * 6.5)
        top = Inches(1 + row * 2)
        
        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.2), Inches(0.4))
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()
        
        header_text = slide.shapes.add_textbox(left, top + Inches(0.05), Inches(6.2), Inches(0.4))
        tf = header_text.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide.shapes.add_textbox(left, top + Inches(0.5), Inches(6.2), Inches(1.3))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
    
    # ========== FEATURE WORKFLOW SLIDES FOR TESTERS ==========
    
    # Feature 1: Patient Registration
    create_feature_workflow_slide(prs, "Patient Registration & ABHA Integration",
        [
            "User opens HealthTrack Pro app/web portal",
            "Clicks 'Register New Patient' button",
            "Enters patient demographics (name, DOB, gender, contact)",
            "System validates mobile number via OTP",
            "Option to link existing ABHA ID or create new",
            "If new ABHA: Aadhaar-based verification flow",
            "System creates patient profile with unique ID",
            "Health history form (optional) appears",
            "Patient receives confirmation SMS/WhatsApp",
            "Profile synced to ABDM health locker"
        ],
        [
            "Verify OTP delivery within 30 seconds",
            "Test invalid mobile number handling",
            "Verify ABHA ID validation against ABDM",
            "Test duplicate patient detection",
            "Verify offline registration queuing",
            "Test multi-language form rendering",
            "Verify data encryption at rest",
            "Test session timeout handling",
            "Verify consent capture compliance",
            "Test accessibility (screen reader)"
        ])
    
    # Feature 2: Telemedicine Consultation
    create_feature_workflow_slide(prs, "Telemedicine Video Consultation",
        [
            "Patient logs in and selects 'Book Consultation'",
            "Chooses specialty and views available doctors",
            "Selects preferred doctor and time slot",
            "Pays consultation fee (UPI/Card/Wallet)",
            "Receives booking confirmation with meeting link",
            "At appointment time, joins video call",
            "Doctor reviews patient history in split-screen",
            "AI assistant suggests relevant questions",
            "Doctor examines patient via video",
            "E-prescription generated and sent to patient",
            "Follow-up reminder scheduled automatically"
        ],
        [
            "Test video quality in low bandwidth (2G/3G)",
            "Verify doctor-patient audio sync",
            "Test screen sharing functionality",
            "Verify payment gateway integration",
            "Test refund flow for cancelled appointments",
            "Verify prescription PDF generation",
            "Test multi-participant call (family member)",
            "Verify recording consent capture",
            "Test call reconnection after drop",
            "Verify AI suggestion accuracy"
        ])
    
    # Feature 3: AI Diagnosis Assistant
    create_feature_workflow_slide(prs, "AI-Powered Diagnosis Assistant",
        [
            "Doctor opens patient consultation screen",
            "Clicks 'AI Assistant' button",
            "Enters/speaks patient symptoms",
            "AI analyzes symptoms against disease database",
            "Displays probable conditions with confidence scores",
            "Shows recommended tests for each condition",
            "Doctor can ask follow-up questions to AI",
            "AI provides differential diagnosis matrix",
            "Doctor selects/modifies final diagnosis",
            "AI suggestions logged for audit trail"
        ],
        [
            "Test symptom recognition accuracy (50+ conditions)",
            "Verify confidence score calculation",
            "Test voice input in multiple languages",
            "Verify AI response time (<3 seconds)",
            "Test offline AI inference capability",
            "Verify medical terminology accuracy",
            "Test edge cases (rare diseases)",
            "Verify audit log completeness",
            "Test AI disclaimer display",
            "Verify doctor override logging"
        ])
    
    # Feature 4: Vitals Monitoring
    create_feature_workflow_slide(prs, "Real-time Vitals Monitoring",
        [
            "Patient pairs wearable device via Bluetooth",
            "System detects device type (smartwatch, BP monitor, etc.)",
            "Initial sync downloads historical data",
            "Real-time vitals display on dashboard",
            "Alert thresholds configured per patient",
            "Abnormal reading triggers notification",
            "Doctor receives alert with patient context",
            "Historical trends displayed in graphs",
            "Export option for vitals report",
            "Integration with consultation notes"
        ],
        [
            "Test Bluetooth pairing with 10+ device brands",
            "Verify data sync accuracy (±1% variance)",
            "Test alert trigger timing (<30 seconds)",
            "Verify offline data buffering",
            "Test battery optimization impact",
            "Verify graph rendering performance",
            "Test concurrent device connections",
            "Verify data privacy (encrypted transmission)",
            "Test alert escalation workflow",
            "Verify export format compliance (FHIR)"
        ])
    
    # Feature 5: E-Prescription
    create_feature_workflow_slide(prs, "Digital Prescription Management",
        [
            "Doctor completes consultation",
            "Opens prescription module",
            "Searches medicine database (generic + brand)",
            "Selects medicines with dosage instructions",
            "System checks drug interactions automatically",
            "Adds special instructions if needed",
            "Digital signature applied",
            "Prescription saved to patient record",
            "PDF/image sent to patient via SMS/WhatsApp",
            "Pharmacy notification for medicine availability"
        ],
        [
            "Test medicine search (1L+ database)",
            "Verify drug interaction alerts",
            "Test dosage calculation accuracy",
            "Verify digital signature validity",
            "Test PDF generation quality",
            "Verify pharmacy API integration",
            "Test prescription edit/cancel flow",
            "Verify refill reminder scheduling",
            "Test insurance pre-authorization",
            "Verify prescription expiry handling"
        ])
    
    # Feature 6: Lab Report Analysis & AI Insights (CORE FEATURE - NOT TEST BOOKING)
    create_feature_workflow_slide(prs, "Lab Report Pull & AI Analysis (Core Feature)",
        [
            "User connects to lab network OR uploads existing report (PDF/image)",
            "System PULLS reports directly from partner diagnostic labs",
            "AI extracts ALL parameters using OCR + NLP (100+ parameters)",
            "Normalizes values across different lab formats & reference ranges",
            "AI analyzes results against healthy ranges for age/gender",
            "Correlates with wearable data (sleep, HRV, activity patterns)",
            "Generates comprehensive health score with trend analysis",
            "Identifies parameters needing attention (high/low flags)",
            "Creates personalized lifestyle recommendations",
            "Tracks ALL parameters over time with visual graphs",
            "Option to share detailed report with consulting doctor"
        ],
        [
            "Test lab network API integration (10+ lab partners)",
            "Test OCR accuracy on 30+ lab report formats",
            "Verify parameter extraction (100+ parameters)",
            "Test value normalization across lab standards",
            "Verify AI analysis accuracy vs medical benchmarks",
            "Test wearable data correlation logic",
            "Verify health score calculation algorithm",
            "Test trend detection over 6+ months of data",
            "Verify recommendation relevance to parameters",
            "Test PDF export of comprehensive report"
        ])
    
    # Feature 7: Wearable Data Integration - Comprehensive
    create_feature_workflow_slide(prs, "Complete Wearable & Health App Integration",
        [
            "Connect Apple Health / Samsung Health / Google Fit",
            "SLEEP DATA:",
            "  - Total sleep duration, time in bed",
            "  - Sleep stages: REM, Deep Sleep, Light Sleep",
            "  - Sleep quality score, interruptions",
            "HEART DATA:",
            "  - Resting heart rate, active heart rate",
            "  - Heart Rate Variability (HRV) - stress indicator",
            "  - Blood oxygen (SpO2), respiratory rate",
            "ACTIVITY DATA:",
            "  - Steps, distance, floors climbed",
            "  - Active calories, total calories",
            "  - Exercise minutes, workout types",
            "ADDITIONAL:",
            "  - Stress levels, mindfulness minutes",
            "  - Menstrual cycle tracking (if applicable)",
            "  - Blood pressure, glucose (connected devices)"
        ],
        [
            "Test Apple HealthKit full API integration",
            "Test Samsung Health SDK all data types",
            "Test Google Fit API comprehensive sync",
            "Verify sleep stage accuracy vs device",
            "Test HRV calculation and trending",
            "Verify SpO2 continuous monitoring",
            "Test historical import (90+ days)",
            "Verify real-time sync (every 15 min)",
            "Test data encryption and privacy",
            "Verify benchmark comparisons"
        ])
    
    # Feature 8: Ayurvedic Lifestyle Recommendations (CLEAR DISCLAIMER)
    create_feature_workflow_slide(prs, "Ayurvedic Lifestyle Suggestions (NOT Medical Advice)",
        [
            "IMPORTANT: All below are LIFESTYLE SUGGESTIONS, not medical recommendations",
            "",
            "AI determines user's Prakriti (Vata/Pitta/Kapha constitution)",
            "Analyzes lab reports + wearable data holistically",
            "",
            "YOGA SUGGESTIONS:",
            "  - Specific asanas based on health parameters",
            "  - Morning/evening routines, duration guidance",
            "",
            "DIET PLAN (Ahara):",
            "  - Personalized Ayurvedic diet (Sattvic/Rajasic/Tamasic)",
            "  - Foods to favor/avoid based on constitution",
            "  - Meal timing aligned with Ayurvedic principles",
            "",
            "DAILY ROUTINE (Dinacharya):",
            "  - Optimal wake/sleep times",
            "  - Meal schedules, exercise windows",
            "  - Self-care practices (Abhyanga, etc.)",
            "",
            "HERBAL SUGGESTIONS:",
            "  - Traditional Ayurvedic herbs and formulations",
            "  - Seasonal guidance (Ritucharya)"
        ],
        [
            "Test Prakriti assessment accuracy",
            "Verify yoga recommendation logic",
            "Test diet plan personalization",
            "Verify Dinacharya scheduling",
            "Test herbal suggestion database",
            "Verify seasonal adjustments",
            "TEST DISCLAIMER VISIBILITY ON ALL SCREENS",
            "Verify 'Consult Doctor' CTA prominent",
            "Test multilingual recommendations",
            "Verify PDF export with disclaimer"
        ])
    
    # Feature 9: Ayurvedic Doctor Platform - Global Vision
    create_feature_workflow_slide(prs, "Ayurvedic Doctor Platform - Taking Ayurveda Global",
        [
            "VERIFIED PRACTITIONERS:",
            "  - BAMS/MD Ayurveda verified credentials",
            "  - Specialization-wise listing (Panchakarma, Rasayana, etc.)",
            "",
            "USER EXPERIENCE:",
            "  - Browse doctors by specialty/language/rating",
            "  - Book video/audio consultation with Vaidya",
            "  - Doctor reviews lab + wearable + AI insights",
            "",
            "TREATMENT PLANS:",
            "  - Personalized Ayurvedic treatment protocols",
            "  - Authentic Ayurvedic medicine prescriptions",
            "  - Integrated Ayurvedic pharmacy network",
            "",
            "GLOBAL VISION:",
            "  - Multi-language support for international users",
            "  - Multi-currency payments (USD, EUR, GBP, etc.)",
            "  - Timezone-aware scheduling",
            "  - Making authentic Ayurveda accessible worldwide",
            "",
            "FOLLOW-UP:",
            "  - Progress tracking and monitoring",
            "  - Doctor-patient messaging"
        ],
        [
            "Test doctor BAMS/MD verification workflow",
            "Verify profile completeness validation",
            "Test specialty search and filters",
            "Verify HD video consultation quality",
            "Test Ayurvedic prescription format",
            "Verify pharmacy network integration",
            "Test international user registration",
            "Verify multi-currency payments",
            "Test timezone scheduling accuracy",
            "Verify doctor rating system"
        ])
    
    # Feature 10: Offline Mode
    create_feature_workflow_slide(prs, "Offline Mode Operations",
        [
            "App detects loss of internet connectivity",
            "Switches to offline mode automatically",
            "All patient data available locally (encrypted)",
            "New registrations queued for sync",
            "Consultations can be recorded offline",
            "Prescriptions generated with offline signature",
            "Vitals continue to be collected",
            "Sync indicator shows pending items",
            "On reconnection, background sync starts",
            "Conflict resolution for simultaneous edits"
        ],
        [
            "Test auto-detection of connectivity loss",
            "Verify local database encryption",
            "Test queue management (1000+ items)",
            "Verify offline consultation recording",
            "Test sync priority ordering",
            "Verify conflict resolution logic",
            "Test partial sync recovery",
            "Verify data integrity post-sync",
            "Test offline mode duration (72+ hours)",
            "Verify battery usage in offline mode"
        ])
    
    # Feature 8: Multi-language Support
    create_feature_workflow_slide(prs, "Multi-Language Interface",
        [
            "User selects preferred language at login",
            "Entire UI renders in selected language",
            "Voice commands accepted in local language",
            "AI responses generated in user language",
            "Prescriptions can be printed in patient's language",
            "SMS/WhatsApp notifications in preferred language",
            "Help documentation available in all languages",
            "Language can be changed mid-session",
            "Medical terms have local equivalents",
            "Right-to-left support for Urdu/Arabic"
        ],
        [
            "Test all 8 supported languages",
            "Verify font rendering quality",
            "Test voice recognition accuracy per language",
            "Verify translation accuracy (medical terms)",
            "Test language switching performance",
            "Verify RTL layout for Urdu",
            "Test SMS character encoding",
            "Verify PDF generation in all languages",
            "Test accessibility with local screen readers",
            "Verify date/number formatting per locale"
        ])
    
    # ========== 3-YEAR ROADMAP ==========
    create_roadmap_slide(prs, "HealthTrack Pro - 3 Year Roadmap",
        [
            ("YEAR 1 (2025)", [
                "Launch lab integration (10+ lab networks)",
                "Apple Health & Samsung Health integration",
                "Prakriti assessment & basic recommendations",
                "Yoga, Diet & Daily Routine suggestions",
                "Onboard 100+ verified Ayurvedic doctors",
                "8 Indian languages supported",
                "1 million user registrations",
                "Mobile apps (iOS + Android)",
                "Clear disclaimer framework",
                "₹15 Cr ARR target"
            ]),
            ("YEAR 2 (2026)", [
                "Pan-India lab network coverage",
                "Google Fit & Fitbit integration",
                "Advanced AI correlation engine",
                "500+ Ayurvedic doctors on platform",
                "International expansion (US, UK, UAE)",
                "Multi-currency payment gateway",
                "15 languages (Hindi, English, Intl)",
                "5 million users globally",
                "Ayurvedic pharmacy network integration",
                "₹50 Cr ARR target"
            ]),
            ("YEAR 3 (2027)", [
                "Global lab network partnerships",
                "All major wearable platforms",
                "Predictive wellness AI",
                "2000+ Ayurvedic doctors worldwide",
                "20+ countries active",
                "Ayurveda certification courses",
                "Research partnerships (Ayurveda validation)",
                "10 million users globally",
                "Enterprise wellness programs",
                "₹100 Cr ARR - Taking Ayurveda Global!"
            ])
        ])
    
    # ========== SLIDE: Complete Feature List ==========
    create_content_slide(prs, "Complete Feature List - Core Capabilities", [
        "LAB REPORT ANALYSIS: Auto-pull from labs, OCR extraction, AI analysis, parameter tracking, trend visualization",
        "WEARABLE INTEGRATION: Apple Health, Samsung Health, Google Fit - sleep stages, HRV, SpO2, steps, stress",
        "AYURVEDIC WELLNESS: Prakriti assessment, personalized yoga, diet plans, daily routine (Dinacharya), herbal suggestions",
        "DOCTOR PLATFORM: Verified Ayurvedic practitioners (BAMS/MD), video consultations, treatment plans, follow-ups",
        "GLOBAL ACCESS: Multi-language (8+ languages), multi-currency payments, timezone scheduling, taking Ayurveda worldwide",
        "AI INSIGHTS: Correlation analysis (lab + wearable data), health scoring, personalized recommendations",
        "HEALTH TRACKING: Parameter trends over time, comparative analysis, benchmark scoring, progress reports",
        "DISCLAIMER COMPLIANCE: All suggestions clearly marked as lifestyle advice, not medical recommendations",
        "DATA PRIVACY: End-to-end encryption, user data ownership, GDPR/data protection compliant",
        "OFFLINE MODE: Full functionality without internet, auto-sync on reconnection"
    ], INFUSE_ORANGE)
    
    # ========== SLIDE: Key Differentiators ==========
    create_content_slide(prs, "Why HealthTrack Pro is Different", [
        "NOT A TEST BOOKING PLATFORM: We don't book tests - we pull and analyze existing reports from any lab",
        "WEARABLE-FIRST APPROACH: Deep integration with Apple Health, Samsung Health - analyzing ALL available data",
        "AYURVEDA + AI FUSION: Ancient wisdom validated by modern data - personalized lifestyle recommendations",
        "CLEAR BOUNDARIES: Lifestyle suggestions only - we always recommend consulting qualified medical professionals",
        "DOCTOR NETWORK: Verified Ayurvedic practitioners for those wanting personalized guidance",
        "GLOBAL MISSION: Taking authentic Ayurveda to the world - not just India",
        "HOLISTIC ANALYSIS: Correlating lab parameters with wearable data for complete health picture",
        "ACTIONABLE INSIGHTS: Not just numbers - specific yoga, diet, and routine suggestions you can follow",
        "TREND TRACKING: See how your health parameters change over time with lifestyle modifications",
        "PRIVACY-FIRST: Your health data belongs to you - full encryption and data ownership"
    ], INFUSE_DEEP_ORANGE)
    
    # ========== SLIDE: API Endpoints for Testing ==========
    create_content_slide(prs, "Key API Endpoints for Testing", [
        "POST /api/auth/login - User authentication with JWT tokens",
        "POST /api/labs/connect - Connect to lab network for report pulling",
        "GET /api/labs/reports/{user_id} - Fetch pulled lab reports",
        "POST /api/labs/analyze - AI analysis of lab parameters",
        "POST /api/wearables/connect - Connect Apple/Samsung/Google Health",
        "GET /api/wearables/data/{user_id} - Fetch all wearable health data",
        "POST /api/health/correlate - Correlate lab + wearable data",
        "GET /api/health/score/{user_id} - Get comprehensive health score",
        "POST /api/ayurveda/prakriti - Assess user's Ayurvedic constitution",
        "GET /api/ayurveda/recommendations/{user_id} - Get lifestyle suggestions",
        "GET /api/doctors/ayurvedic - List verified Ayurvedic doctors",
        "POST /api/consultations/book - Book doctor consultation"
    ], INFUSE_DEEP_ORANGE)
    
    # ========== SLIDE: Team ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_logo(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Leadership Team"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = INFUSE_GOLD
    p.alignment = PP_ALIGN.CENTER
    
    team = [
        ("Rohini Koul", "CEO", "20 years in Academia\nMSc Chemistry", INFUSE_GOLD),
        ("Chief Growth Officer", "Founder", "30 years in Technology\nMBA, University of Arizona", INFUSE_ORANGE),
        ("Chief Technology Officer", "CTO", "24 years in SaaS/PaaS\nMBA, University of Pennsylvania", INFUSE_DEEP_ORANGE)
    ]
    
    for i, (name, role, bio, color) in enumerate(team):
        left = Inches(0.8 + i * 4.2)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(3.8), Inches(5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x2D, 0x33, 0x3F)
        card.line.fill.background()
        
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.9), Inches(1.7), Inches(2), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        
        badge_text = slide.shapes.add_textbox(left + Inches(0.9), Inches(1.7), Inches(2), Inches(0.4))
        tf = badge_text.text_frame
        p = tf.paragraphs[0]
        p.text = role
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE if color != INFUSE_GOLD else DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
        
        name_box = slide.shapes.add_textbox(left, Inches(2.3), Inches(3.8), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        bio_box = slide.shapes.add_textbox(left + Inches(0.2), Inches(2.9), Inches(3.4), Inches(2))
        tf = bio_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bio
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    # ========== SLIDE: Contact ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_logo(slide)
    add_accent_bar(slide, INFUSE_ORANGE, 3.2)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.33), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Let's Transform Healthcare Together"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    contact_info = [
        "info@infuse.net.in",
        "+91-9599960663",
        "www.infuse.net.in"
    ]
    
    for i, info in enumerate(contact_info):
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2 + i * 0.6), Inches(12.33), Inches(0.5))
        tf = info_box.text_frame
        p = tf.paragraphs[0]
        p.text = info
        p.font.size = Pt(22)
        p.font.color.rgb = INFUSE_GOLD if i == 0 else WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Save
    output_path = "/app/frontend/public/HealthTrack_Pro_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"HealthTrack Pro Pitch Deck saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_healthtrack_pitch_deck()
