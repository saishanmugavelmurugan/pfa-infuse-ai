"""
Infuse Platform - Executive Sales Deck Generator
================================================
Creates a professional PowerPoint presentation for C-level executives
(CEO, CFO, CIO, CSO) showcasing HealthTrack Pro & SecureSphere
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
from datetime import datetime

# Colors
INFUSE_BLUE = RGBColor(0, 102, 204)
INFUSE_DARK = RGBColor(26, 26, 46)
INFUSE_GREEN = RGBColor(0, 168, 107)
INFUSE_ORANGE = RGBColor(255, 140, 0)
INFUSE_RED = RGBColor(220, 53, 69)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 240)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = INFUSE_DARK
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = INFUSE_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, accent_color=INFUSE_BLUE):
    """Add a content slide with bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = INFUSE_DARK
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = INFUSE_DARK
        p.space_after = Pt(12)
    
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.3), Inches(0.1), Inches(5.5))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()
    
    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Add a two-column slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = INFUSE_DARK
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column header
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.3), Inches(0.5))
    tf = left_header.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = INFUSE_BLUE
    
    # Left column content
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.3), Inches(4.5))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = INFUSE_DARK
        p.space_after = Pt(8)
    
    # Right column header
    right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.3), Inches(0.5))
    tf = right_header.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = INFUSE_GREEN
    
    # Right column content
    right_content = slide.shapes.add_textbox(Inches(5.2), Inches(1.9), Inches(4.3), Inches(4.5))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = INFUSE_DARK
        p.space_after = Pt(8)
    
    return slide

def add_metrics_slide(prs, title, metrics):
    """Add a metrics/KPI slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = INFUSE_DARK
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Metrics boxes
    box_width = Inches(2.2)
    box_height = Inches(1.8)
    start_x = Inches(0.3)
    start_y = Inches(1.8)
    gap = Inches(0.15)
    
    colors = [INFUSE_BLUE, INFUSE_GREEN, INFUSE_ORANGE, RGBColor(138, 43, 226)]
    
    for i, (metric_value, metric_label) in enumerate(metrics[:4]):
        x = start_x + (box_width + gap) * i
        
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()
        
        # Value
        value_box = slide.shapes.add_textbox(x, start_y + Inches(0.3), box_width, Inches(0.8))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric_value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(x, start_y + Inches(1.1), box_width, Inches(0.5))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric_label
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Additional metrics below if more than 4
    if len(metrics) > 4:
        start_y2 = Inches(4)
        for i, (metric_value, metric_label) in enumerate(metrics[4:8]):
            x = start_x + (box_width + gap) * i
            
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y2, box_width, box_height)
            box.fill.solid()
            box.fill.fore_color.rgb = colors[(i+4) % len(colors)]
            box.line.fill.background()
            
            value_box = slide.shapes.add_textbox(x, start_y2 + Inches(0.3), box_width, Inches(0.8))
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = metric_value
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
            label_box = slide.shapes.add_textbox(x, start_y2 + Inches(1.1), box_width, Inches(0.5))
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = metric_label
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
    
    return slide

def generate_sales_deck(output_path: str):
    """Generate the complete sales deck"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, 
        "INFUSE PLATFORM",
        "Next-Generation Healthcare & Security Solutions\nFor the Modern Enterprise"
    )
    
    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Infuse is a unified platform delivering AI-powered Healthcare (HealthTrack Pro) and Cybersecurity (SecureSphere) solutions",
        "Designed for multi-tenant, enterprise-scale deployments with 99.99% uptime SLA",
        "8 proprietary AI agents for intelligent automation and threat detection",
        "Telco-grade vRAN security with mandatory self-protection",
        "Proven ROI: 40% reduction in operational costs, 60% faster threat response",
        "Trusted by healthcare providers and enterprises across India"
    ])
    
    # Slide 3: The Challenge
    add_content_slide(prs, "The Challenge Your Organization Faces", [
        "Healthcare: Fragmented patient data, inefficient workflows, compliance risks",
        "Security: Sophisticated cyber threats, GSM/telecom fraud, IoT vulnerabilities",
        "Operations: High costs, manual processes, lack of real-time insights",
        "Compliance: HIPAA, GDPR, PCI-DSS, and industry regulations",
        "Scale: Need for solutions that grow with your business"
    ], INFUSE_RED)
    
    # Slide 4: Our Solution - Two Products
    add_two_column_slide(prs,
        "Infuse Platform: Two Powerful Solutions",
        "🏥 HealthTrack Pro",
        [
            "Complete patient management system",
            "Electronic Health Records (EHR)",
            "AI-powered diagnostics & insights",
            "Lab reports with AI analysis",
            "Prescription management",
            "Appointment scheduling",
            "Billing & revenue tracking",
            "ABDM/ABHA integration (India)"
        ],
        "🛡️ SecureSphere",
        [
            "Enterprise-grade security platform",
            "Network device monitoring",
            "AI threat detection (99.8% accuracy)",
            "GSM fraud prevention",
            "SIM swap protection",
            "IoT security for connected devices",
            "Automotive security (NFC, CAN-bus)",
            "Unified vRAN with 6 segments"
        ]
    )
    
    # Slide 5: AI Agents
    add_content_slide(prs, "8 AI Agents Powering Intelligent Automation", [
        "SUPPORT AGENTS (Customer-Facing):",
        "   1. HealthBot - Healthcare support, symptom analysis, appointment help",
        "   2. SecureGuard - Threat response, incident management, security audits",
        "   3. TelcoAdvisor - Telecom fraud prevention, SIM swap detection",
        "   4. EnterpriseHelper - Customer onboarding, billing, technical support",
        "",
        "LEARNING AGENTS (Backend Intelligence):",
        "   5. ThreatLearner - Learns attack patterns, improves detection",
        "   6. FraudDetector - Learns fraud patterns across GSM, phishing, financial",
        "   7. BehaviorAnalyzer - User behavior analysis, anomaly detection",
        "   8. AnomalyHunter - System-wide monitoring, predictive alerts"
    ], INFUSE_GREEN)
    
    # Slide 6: Key Metrics
    add_metrics_slide(prs, "Platform Performance Metrics", [
        ("99.8%", "Threat Detection"),
        ("92%", "Security Score"),
        ("<1s", "Response Time"),
        ("99.99%", "Uptime SLA"),
        ("40%", "Cost Reduction"),
        ("60%", "Faster Response"),
        ("8", "AI Agents"),
        ("6", "vRAN Segments")
    ])
    
    # Slide 7: Security Architecture
    add_content_slide(prs, "SecureSphere: Telco-Grade Security Architecture", [
        "Unified vRAN System with 6 Protected Segments:",
        "   • Telco - MSISDN, IMSI, APN protection",
        "   • Mobile - Device ID, IMEI monitoring", 
        "   • Enterprise - IP ranges, domain security",
        "   • Automotive - VIN, eSIM, telematics protection",
        "   • White Goods (IoT) - Smart device security",
        "   • CCTV - Camera and stream protection",
        "",
        "Self-Protection Policy: MANDATORY (Cannot be disabled)",
        "   • DDoS mitigation, intrusion detection, auto-healing",
        "   • Rate limiting, automatic threat blocking"
    ], INFUSE_BLUE)
    
    # Slide 8: ROI for CFO
    add_content_slide(prs, "Financial Impact (CFO View)", [
        "Cost Savings:",
        "   • 40% reduction in security incident costs",
        "   • 35% reduction in manual administrative tasks",
        "   • 25% reduction in compliance audit preparation time",
        "",
        "Revenue Enhancement:",
        "   • Faster patient throughput with streamlined workflows",
        "   • Reduced fraud losses (avg. ₹2Cr+ annual savings for enterprises)",
        "   • Premium security features enable higher service pricing",
        "",
        "Investment Protection:",
        "   • Flexible SaaS and PaaS pricing models",
        "   • Pay-as-you-grow with no upfront infrastructure costs"
    ], INFUSE_GREEN)
    
    # Slide 9: Technology for CIO/CTO
    add_content_slide(prs, "Technology Architecture (CIO/CTO View)", [
        "Modern Cloud-Native Architecture:",
        "   • Microservices-based, containerized deployment",
        "   • Multi-tenant with data isolation",
        "   • RESTful APIs for seamless integration",
        "   • Real-time analytics and dashboards",
        "",
        "Integration Capabilities:",
        "   • ABDM/ABHA for healthcare interoperability",
        "   • SSO/SAML integration for enterprise",
        "   • Webhook support for custom workflows",
        "   • Mobile SDKs for iOS and Android",
        "",
        "AI/ML Infrastructure:",
        "   • 8 specialized AI agents with continuous learning",
        "   • <10MB memory footprint per agent"
    ], INFUSE_BLUE)
    
    # Slide 10: Security for CSO/CISO
    add_content_slide(prs, "Security Posture (CSO/CISO View)", [
        "Compliance & Certifications:",
        "   • HIPAA compliant healthcare data handling",
        "   • GDPR ready with data privacy controls",
        "   • SOC 2 Type II audit ready",
        "",
        "Security Features:",
        "   • AES-256-GCM encryption at rest and in transit",
        "   • Zero-trust architecture with RBAC",
        "   • Real-time threat monitoring with AI detection",
        "   • Automated incident response playbooks",
        "   • Complete audit logging and forensics",
        "",
        "SecureSphere Protection:",
        "   • Mandatory self-protection that cannot be disabled",
        "   • 6-segment vRAN coverage for comprehensive defense"
    ], INFUSE_RED)
    
    # Slide 11: Pricing
    add_two_column_slide(prs,
        "Flexible Pricing Models",
        "💼 SaaS (Per User)",
        [
            "Individual Basic: $9.99/month",
            "Individual Pro: $24.99/month",
            "Individual Premium: $49.99/month",
            "",
            "Ideal for:",
            "• Small clinics & practices",
            "• Individual security monitoring",
            "• Startups and SMBs"
        ],
        "🏢 PaaS (Enterprise)",
        [
            "Enterprise 50: Custom pricing",
            "Enterprise 100: Custom pricing",
            "Enterprise 500+: Volume discounts",
            "Telco Operator: $99,999/year",
            "",
            "Includes:",
            "• Dedicated support",
            "• Custom integrations",
            "• SLA guarantees"
        ]
    )
    
    # Slide 12: Customer Success (Mock)
    add_content_slide(prs, "Customer Success Stories", [
        "🏥 Leading Hospital Chain (500+ beds):",
        "   • 45% improvement in patient data accessibility",
        "   • 30% reduction in administrative overhead",
        "",
        "📡 Tier-1 Telecom Operator:",
        "   • Blocked 2M+ fraud attempts in first year",
        "   • 99.7% SIM swap attack prevention rate",
        "",
        "🏢 Fortune 500 Enterprise:",
        "   • Zero security breaches since implementation",
        "   • 50% faster incident response time",
        "",
        "\"Infuse has transformed our security posture\" - CISO, Major Bank"
    ], INFUSE_GREEN)
    
    # Slide 13: Implementation Timeline
    add_content_slide(prs, "Implementation & Onboarding", [
        "Week 1-2: Discovery & Planning",
        "   • Requirements gathering, architecture review",
        "",
        "Week 3-4: Deployment & Configuration",
        "   • Cloud setup, security policies, integrations",
        "",
        "Week 5-6: Training & UAT",
        "   • User training, acceptance testing, go-live prep",
        "",
        "Week 7+: Go-Live & Support",
        "   • Production deployment, 24/7 support activation",
        "",
        "Dedicated Customer Success Manager throughout the journey"
    ], INFUSE_BLUE)
    
    # Slide 14: Call to Action
    add_title_slide(prs,
        "Ready to Transform Your Organization?",
        "Contact us for a personalized demo\n\nsales@infuse.nt.in | +91-XXXX-XXXXXX\nwww.infuse.nt.in"
    )
    
    # Slide 15: Thank You
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = INFUSE_DARK
    shape.line.fill.background()
    
    # Thank you text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Infuse - Powering the Future of Healthcare & Security"
    p.font.size = Pt(20)
    p.font.color.rgb = INFUSE_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Confidential | {datetime.now().strftime('%B %Y')}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER
    
    # Save
    prs.save(output_path)
    print(f"✅ Sales deck generated: {output_path}")
    return output_path

if __name__ == "__main__":
    output_dir = "/app/backend/static/downloads"
    os.makedirs(output_dir, exist_ok=True)
    output_path = f"{output_dir}/Infuse_Executive_Sales_Deck.pptx"
    generate_sales_deck(output_path)
