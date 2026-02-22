from fastapi import APIRouter
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io

router = APIRouter()

@router.get("/generate-pitch-deck")
async def generate_pitch_deck():
    """Generate and download Infuse-ai pitch deck as PowerPoint"""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    ORANGE = RGBColor(249, 115, 22)
    DARK_ORANGE = RGBColor(234, 88, 12)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(55, 65, 81)
    LIGHT_GRAY = RGBColor(243, 244, 246)
    GREEN = RGBColor(16, 185, 129)
    RED = RGBColor(239, 68, 68)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = ORANGE
    
    # Add title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Infuse-ai"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Enterprise Healthcare & Security Solutions"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = WHITE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add description
    desc_box = slide1.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
    desc_frame = desc_box.text_frame
    desc_frame.text = "Transforming Healthcare Management and IoT Security with AI-Powered Intelligence"
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(18)
    desc_para.font.color.rgb = WHITE
    desc_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: The Problem
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    tf = title.text_frame
    tf.text = "The Problem"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GRAY
    
    # Subtitle
    subtitle = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.4))
    tf2 = subtitle.text_frame
    tf2.text = "Critical Gaps in Healthcare & Security"
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(20)
    p2.font.color.rgb = GRAY
    
    # Healthcare Crisis
    health_title = slide2.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4.5), Inches(0.5))
    health_title.text_frame.text = "Healthcare Crisis"
    health_title.text_frame.paragraphs[0].font.size = Pt(24)
    health_title.text_frame.paragraphs[0].font.bold = True
    health_title.text_frame.paragraphs[0].font.color.rgb = RED
    
    health_problems = [
        "• Fragmented Records: Patient data scattered across providers",
        "• Expensive Care: Specialist consultations cost $199-$999",
        "• Poor Insurance: Manual claims, no real-time verification",
        "• Urban-Centric: Limited reach in rural areas"
    ]
    
    health_box = slide2.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(4.5), Inches(2.5))
    tf3 = health_box.text_frame
    for problem in health_problems:
        p = tf3.add_paragraph()
        p.text = problem
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.space_after = Pt(10)
    
    # Security Challenges
    sec_title = slide2.shapes.add_textbox(Inches(5.5), Inches(2.0), Inches(4.5), Inches(0.5))
    sec_title.text_frame.text = "Security Challenges"
    sec_title.text_frame.paragraphs[0].font.size = Pt(24)
    sec_title.text_frame.paragraphs[0].font.bold = True
    sec_title.text_frame.paragraphs[0].font.color.rgb = ORANGE
    
    sec_problems = [
        "• Vendor Lock-in: Telcos trapped with RAN vendors",
        "• Detection-Only: Systems alert but don't mitigate",
        "• Revenue Losses: $10M+/year in fraud & breaches",
        "• Complex Ops: Multiple systems, high overhead"
    ]
    
    sec_box = slide2.shapes.add_textbox(Inches(5.5), Inches(2.6), Inches(4.5), Inches(2.5))
    tf4 = sec_box.text_frame
    for problem in sec_problems:
        p = tf4.add_paragraph()
        p.text = problem
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.space_after = Pt(10)
    
    # Slide 3: Our Solution
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title3.text_frame.text = "Our Solution"
    title3.text_frame.paragraphs[0].font.size = Pt(44)
    title3.text_frame.paragraphs[0].font.bold = True
    title3.text_frame.paragraphs[0].font.color.rgb = GRAY
    
    subtitle3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.4))
    subtitle3.text_frame.text = "Two Revolutionary Products"
    subtitle3.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle3.text_frame.paragraphs[0].font.color.rgb = GRAY
    
    # HealthTrack Pro
    health_title3 = slide3.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(4.1), Inches(0.5))
    health_title3.text_frame.text = "HealthTrack Pro"
    health_title3.text_frame.paragraphs[0].font.size = Pt(28)
    health_title3.text_frame.paragraphs[0].font.bold = True
    health_title3.text_frame.paragraphs[0].font.color.rgb = DARK_ORANGE
    
    health_sub = slide3.shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(4.1), Inches(0.3))
    health_sub.text_frame.text = "Unified AI-Powered Health Management"
    health_sub.text_frame.paragraphs[0].font.size = Pt(14)
    health_sub.text_frame.paragraphs[0].font.color.rgb = GRAY
    
    health_features = [
        "✓ Unified Health Vault: All medical records",
        "✓ Smart Watch Integration: Real-time monitoring",
        "✓ AI Analytics: 80% prediction accuracy",
        "✓ Doctor Network: 10,000+ verified providers",
        "✓ Insurance Integration: One-click claims"
    ]
    
    health_feat_box = slide3.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(4.1), Inches(1.8))
    tf5 = health_feat_box.text_frame
    for feature in health_features:
        p = tf5.add_paragraph()
        p.text = feature
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.space_after = Pt(5)
    
    # SecureIT+IoT
    sec_title3 = slide3.shapes.add_textbox(Inches(5.7), Inches(2.2), Inches(4.1), Inches(0.5))
    sec_title3.text_frame.text = "SecureIT+IoT"
    sec_title3.text_frame.paragraphs[0].font.size = Pt(28)
    sec_title3.text_frame.paragraphs[0].font.bold = True
    sec_title3.text_frame.paragraphs[0].font.color.rgb = GRAY
    
    sec_sub3 = slide3.shapes.add_textbox(Inches(5.7), Inches(2.8), Inches(4.1), Inches(0.3))
    sec_sub3.text_frame.text = "Enterprise Telco-Grade Security"
    sec_sub3.text_frame.paragraphs[0].font.size = Pt(14)
    sec_sub3.text_frame.paragraphs[0].font.color.rgb = GRAY
    
    sec_features = [
        "✓ RAN Independent: Works with ANY vendor",
        "✓ Real-Time Enforcement: <100ms mitigation",
        "✓ Threat Detection: 99.8% ML accuracy",
        "✓ Operational Excellence: 40% NOC reduction",
        "✓ Global Scale: 10M+ devices, 99.999% uptime"
    ]
    
    sec_feat_box = slide3.shapes.add_textbox(Inches(5.7), Inches(3.3), Inches(4.1), Inches(1.8))
    tf6 = sec_feat_box.text_frame
    for feature in sec_features:
        p = tf6.add_paragraph()
        p.text = feature
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.space_after = Pt(5)
    
    # Slide 4: Market Opportunity
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title4.text_frame.text = "Market Opportunity"
    title4.text_frame.paragraphs[0].font.size = Pt(44)
    title4.text_frame.paragraphs[0].font.bold = True
    title4.text_frame.paragraphs[0].font.color.rgb = GRAY
    
    subtitle4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.4))
    subtitle4.text_frame.text = "Massive TAM with Proven Revenue Model"
    subtitle4.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle4.text_frame.paragraphs[0].font.color.rgb = GRAY
    
    # Market boxes
    market1 = slide4.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(4.1), Inches(0.3))
    market1.text_frame.text = "Digital Health Market"
    market1.text_frame.paragraphs[0].font.size = Pt(16)
    market1.text_frame.paragraphs[0].font.color.rgb = ORANGE
    
    market1_val = slide4.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(4.1), Inches(0.6))
    market1_val.text_frame.text = "$639B"
    market1_val.text_frame.paragraphs[0].font.size = Pt(48)
    market1_val.text_frame.paragraphs[0].font.bold = True
    market1_val.text_frame.paragraphs[0].font.color.rgb = ORANGE
    
    market1_desc = slide4.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(4.1), Inches(0.2))
    market1_desc.text_frame.text = "by 2026 (CAGR 27.7%)"
    market1_desc.text_frame.paragraphs[0].font.size = Pt(12)
    market1_desc.text_frame.paragraphs[0].font.color.rgb = GRAY
    
    market2 = slide4.shapes.add_textbox(Inches(5.7), Inches(2.2), Inches(4.1), Inches(0.3))
    market2.text_frame.text = "IoT Security Market"
    market2.text_frame.paragraphs[0].font.size = Pt(16)
    market2.text_frame.paragraphs[0].font.color.rgb = GRAY
    
    market2_val = slide4.shapes.add_textbox(Inches(5.7), Inches(2.6), Inches(4.1), Inches(0.6))
    market2_val.text_frame.text = "$36.6B"
    market2_val.text_frame.paragraphs[0].font.size = Pt(48)
    market2_val.text_frame.paragraphs[0].font.bold = True
    market2_val.text_frame.paragraphs[0].font.color.rgb = GRAY
    
    market2_desc = slide4.shapes.add_textbox(Inches(5.7), Inches(3.2), Inches(4.1), Inches(0.2))
    market2_desc.text_frame.text = "by 2027 (CAGR 21.2%)"
    market2_desc.text_frame.paragraphs[0].font.size = Pt(12)
    market2_desc.text_frame.paragraphs[0].font.color.rgb = GRAY
    
    # Revenue model
    rev_title = slide4.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(0.4))
    rev_title.text_frame.text = "Revenue Model"
    rev_title.text_frame.paragraphs[0].font.size = Pt(20)
    rev_title.text_frame.paragraphs[0].font.bold = True
    rev_title.text_frame.paragraphs[0].font.color.rgb = GRAY
    
    rev_text = slide4.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
    tf7 = rev_text.text_frame
    
    p1 = tf7.paragraphs[0]
    p1.text = "HealthTrack Pro: Free ($0) | Premium ($99/mo) | Enterprise ($299/mo)"
    p1.font.size = Pt(14)
    p1.font.color.rgb = GRAY
    
    p2 = tf7.add_paragraph()
    p2.text = "SecureIT+IoT: 30-Day Trial (Free) | Starter ($9,600/mo) | Pro ($24,000/mo) | Enterprise ($60,000+/mo)"
    p2.font.size = Pt(14)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(10)
    
    # Slide 5: Financial Projections
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title5.text_frame.text = "Financial Projections"
    title5.text_frame.paragraphs[0].font.size = Pt(44)
    title5.text_frame.paragraphs[0].font.bold = True
    title5.text_frame.paragraphs[0].font.color.rgb = GRAY
    
    subtitle5 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.4))
    subtitle5.text_frame.text = "3-Year Revenue Growth with 300% Margin"
    subtitle5.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle5.text_frame.paragraphs[0].font.color.rgb = GRAY
    
    # Year 1
    year1_title = slide5.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(2.6), Inches(0.3))
    year1_title.text_frame.text = "Year 1 (Conservative)"
    year1_title.text_frame.paragraphs[0].font.size = Pt(14)
    year1_title.text_frame.paragraphs[0].font.color.rgb = GREEN
    
    year1_val = slide5.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(2.6), Inches(0.5))
    year1_val.text_frame.text = "$6.44M"
    year1_val.text_frame.paragraphs[0].font.size = Pt(42)
    year1_val.text_frame.paragraphs[0].font.bold = True
    year1_val.text_frame.paragraphs[0].font.color.rgb = GREEN
    
    # Year 2
    year2_title = slide5.shapes.add_textbox(Inches(3.75), Inches(2.2), Inches(2.6), Inches(0.3))
    year2_title.text_frame.text = "Year 2 (Target)"
    year2_title.text_frame.paragraphs[0].font.size = Pt(14)
    year2_title.text_frame.paragraphs[0].font.color.rgb = ORANGE
    
    year2_val = slide5.shapes.add_textbox(Inches(3.75), Inches(2.6), Inches(2.6), Inches(0.5))
    year2_val.text_frame.text = "$23M"
    year2_val.text_frame.paragraphs[0].font.size = Pt(42)
    year2_val.text_frame.paragraphs[0].font.bold = True
    year2_val.text_frame.paragraphs[0].font.color.rgb = ORANGE
    
    # Year 3
    year3_title = slide5.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(2.6), Inches(0.3))
    year3_title.text_frame.text = "Year 3 (Aggressive)"
    year3_title.text_frame.paragraphs[0].font.size = Pt(14)
    year3_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59, 130, 246)
    
    year3_val = slide5.shapes.add_textbox(Inches(7.0), Inches(2.6), Inches(2.6), Inches(0.5))
    year3_val.text_frame.text = "$52.56M"
    year3_val.text_frame.paragraphs[0].font.size = Pt(42)
    year3_val.text_frame.paragraphs[0].font.bold = True
    year3_val.text_frame.paragraphs[0].font.color.rgb = RGBColor(59, 130, 246)
    
    # Slide 6: Call to Action
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    slide6.background.fill.solid()
    slide6.background.fill.fore_color.rgb = ORANGE
    
    cta_title = slide6.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(1))
    cta_title.text_frame.text = "Ready to Get Started?"
    cta_title.text_frame.paragraphs[0].font.size = Pt(48)
    cta_title.text_frame.paragraphs[0].font.bold = True
    cta_title.text_frame.paragraphs[0].font.color.rgb = WHITE
    cta_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    cta_desc = slide6.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.5))
    cta_desc.text_frame.text = "Contact us today to transform your business"
    cta_desc.text_frame.paragraphs[0].font.size = Pt(20)
    cta_desc.text_frame.paragraphs[0].font.color.rgb = WHITE
    cta_desc.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    contact = slide6.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(1.5))
    tf8 = contact.text_frame
    tf8.paragraphs[0].text = "Phone: +1 (555) 123-4567"
    tf8.paragraphs[0].font.size = Pt(16)
    tf8.paragraphs[0].font.color.rgb = WHITE
    tf8.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    p1 = tf8.add_paragraph()
    p1.text = "Email: contact@infuse.net.in"
    p1.font.size = Pt(16)
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf8.add_paragraph()
    p2.text = "Website: www.infuse.net.in"
    p2.font.size = Pt(16)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    # Save to BytesIO
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    # Return as downloadable file
    return StreamingResponse(
        pptx_io,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": "attachment; filename=Infuse-ai-Pitch-Deck.pptx"
        }
    )
